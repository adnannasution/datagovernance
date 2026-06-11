from pptx import Presentation
from pptx.util import Pt
import re

TAG_PATTERN = re.compile(r'\b[0-9][A-Z]{1,2}-[A-Z]{1,3}-[0-9]{3,4}[A-Z]?\b')

def detect_tag_numbers(text: str) -> list[str]:
    return list(set(TAG_PATTERN.findall(text.upper())))

def extract_slide_content(slide) -> tuple[str, str, list[str]]:
    """
    Return (title, body_text, table_rows_serialized)
    """
    title = ""
    body_parts = []
    table_rows = []

    for shape in slide.shapes:
        # Judul slide
        if shape.shape_type == 13:  # picture
            continue

        if hasattr(shape, "text") and shape.text.strip():
            if shape.shape_id == 1 or (hasattr(shape, "name") and
               ("title" in shape.name.lower() or "judul" in shape.name.lower())):
                title = shape.text.strip()
            else:
                body_parts.append(shape.text.strip())

        # Tabel dalam slide
        if shape.has_table:
            tbl = shape.table
            headers = []
            for r_idx, row in enumerate(tbl.rows):
                cells = [cell.text.strip() for cell in row.cells]
                if r_idx == 0:
                    headers = cells
                else:
                    if not any(cells):
                        continue
                    parts = []
                    for h, c in zip(headers, cells):
                        if c:
                            parts.append(f"{h}: {c}")
                    if parts:
                        table_rows.append(" | ".join(parts))

    return title, "\n".join(body_parts), table_rows

def process_pptx(doc_id: int, file_path: str) -> dict:
    prs = Presentation(file_path)
    chunks = []
    table_rows_out = []
    chunk_index = 0
    all_tags = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title, body, tbl_rows = extract_slide_content(slide)

        # Notes
        notes_text = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                notes_text = f"\nCatatan: {notes}"

        # Gabungkan semua konten slide jadi satu chunk
        content_parts = []
        if title:
            content_parts.append(f"[Slide {slide_num}: {title}]")
        else:
            content_parts.append(f"[Slide {slide_num}]")

        if body:
            content_parts.append(body)

        if tbl_rows:
            content_parts.append("[Tabel dalam slide]")
            content_parts.extend(tbl_rows)

            # Simpan tabel slide ke table_rows juga
            for r_idx, row_text in enumerate(tbl_rows):
                tags = detect_tag_numbers(row_text)
                tag_num = tags[0] if tags else None
                table_rows_out.append({
                    "doc_id": doc_id,
                    "halaman": slide_num,
                    "sheet_name": None,
                    "row_index": r_idx,
                    "row_data": {"slide": slide_num, "content": row_text},
                    "tag_number": tag_num
                })

        if notes_text:
            content_parts.append(notes_text)

        full_content = "\n".join(content_parts)

        if full_content.strip():
            chunks.append({
                "doc_id": doc_id,
                "chunk_index": chunk_index,
                "halaman": slide_num,
                "slide_number": slide_num,
                "slide_title": title or f"Slide {slide_num}",
                "sheet_name": None,
                "content": full_content,
                "embedding": None
            })
            chunk_index += 1
            all_tags.extend(detect_tag_numbers(full_content))

    return {
        "chunks": chunks,
        "table_rows": table_rows_out,
        "auto_tags": list(set(all_tags)),
        "total_pages": len(prs.slides)
    }
